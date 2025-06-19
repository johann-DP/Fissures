import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

def create_ppt_from_results(pipeline_results, results_folder, ppt_path):
    """
    Respecte l'ordre exact:
      1) evolution_inclinaison.png
      2) explanatory_evolution_Aggregated_Absolute.png
      3) Les figures model_... en sous-ordre:
         Aggregated_Absolute, Aggregated_Delta,
         Hourly_Absolute, Hourly_Delta,
         AllSimple, AllPenalized
      4) Les figures models_comparison_... dans le même sous-ordre
      5) Tout le reste en ordre alphabétique.
    """

    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    def add_centered_picture(slide, img_path, width_in_inches=8):
        pic = slide.shapes.add_picture(img_path, 0, 0, width=Inches(width_in_inches))
        # Centrage horizontal
        pic.left = (slide_width - pic.width) // 2
        # Centrage vertical
        pic.top = (slide_height - pic.height) // 2

    # Crée la diapo titre
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Résultats de la modélisation de l'inclinaison du mur"
    slide.placeholders[1].text = "Synthèse des résultats et figures générées"

    # Fonction utilitaire pour insérer un PNG en slide (sans titre), centré
    def insert_png(img_path):
        slide_layout = prs.slide_layouts[5]
        new_slide = prs.slides.add_slide(slide_layout)
        if new_slide.shapes.title:
            new_slide.shapes.title.text = ""
        add_centered_picture(new_slide, img_path, width_in_inches=8)

    # 1) evolution_inclinaison.png
    evol_path = os.path.join(results_folder, "evolution_inclinaison.png")
    if os.path.isfile(evol_path):
        insert_png(evol_path)

    # 2) explanatory_evolution_Aggregated_Absolute.png
    expl_path = os.path.join(results_folder, "explanatory_evolution_Aggregated_Absolute.png")
    if os.path.isfile(expl_path):
        insert_png(expl_path)

    # 3) model_... dans l'ordre strict
    #    On recherche tous les .png qui commencent par "model_"
    #    On leur attribue un 'rank' selon la sous-chaîne (Aggregated_Absolute, etc.)
    model_suborder = [
        "Aggregated_Absolute",
        "Aggregated_Delta",
        "Hourly_Absolute",
        "Hourly_Delta",
        "AllSimple",
        "AllPenalized",
    ]

    # # Récupérer tous les fichiers "model_*.png"
    # model_files = []
    # for fn in os.listdir(results_folder):
    #     if fn.lower().endswith(".png") and fn.startswith("model_"):
    #         fullp = os.path.join(results_folder, fn)
    #         # Déterminer le rank
    #         # On prend le premier matching
    #         assigned_rank = None
    #         for i, token in enumerate(model_suborder):
    #             if token in fn:
    #                 assigned_rank = i
    #                 break
    #         # Si pas trouvé, on lui donne un rang après tout le monde
    #         if assigned_rank is None:
    #             assigned_rank = 999
    #         model_files.append((assigned_rank, fn, fullp))
    #
    # # Tri : d'abord par rank, puis alpha sur fn
    # model_files.sort(key=lambda x: (x[0], x[1]))
    # # On insère
    # for (rank, fn, fullp) in model_files:
    #     insert_png(fullp)
    #
    # # 4) models_comparison_... dans le même ordre
    # #    On recherche tous les .png qui commencent par "models_comparison_"
    # comparison_files = []
    # for fn in os.listdir(results_folder):
    #     if fn.lower().endswith(".png") and fn.startswith("models_comparison_"):
    #         fullp = os.path.join(results_folder, fn)
    #         # Déterminer le rank
    #         assigned_rank = None
    #         for i, token in enumerate(model_suborder):
    #             if token in fn:
    #                 assigned_rank = i
    #                 break
    #         if assigned_rank is None:
    #             assigned_rank = 999
    #         comparison_files.append((assigned_rank, fn, fullp))
    #
    # # Tri : d'abord par rank, puis alpha
    # comparison_files.sort(key=lambda x: (x[0], x[1]))
    # for (rank, fn, fullp) in comparison_files:
    #     insert_png(fullp)

    # 5) Tout le reste (les PNG non déjà insérés) par ordre alphabétique
    inserted_files = set()  # Contient le chemin complet

    # On refait la liste des images déjà utilisées
    # (1) evolution_inclinaison
    if os.path.isfile(evol_path):
        inserted_files.add(evol_path)
    # (2) explanatory_evolution
    if os.path.isfile(expl_path):
        inserted_files.add(expl_path)
    # # (3) model_...
    # for _, fn, fullp in model_files:
    #     inserted_files.add(fullp)
    # # (4) models_comparison_...
    # for _, fn, fullp in comparison_files:
    #     inserted_files.add(fullp)

    # Parcours final du dossier
    leftover_pngs = []
    for fn in os.listdir(results_folder):
        if fn.lower().endswith(".png"):
            fullp = os.path.join(results_folder, fn)
            if fullp not in inserted_files:
                leftover_pngs.append(fn)
    leftover_pngs.sort()  # tri alpha
    for fn in leftover_pngs:
        fullp = os.path.join(results_folder, fn)
        insert_png(fullp)

    # On sauvegarde le PPT
    prs.save(ppt_path)
